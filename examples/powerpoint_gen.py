from pptx import Presentation

from pptx.util import Inches



# Create a presentation object

prs = Presentation()



# Define slide layout

title_slide_layout = prs.slide_layouts[0] # 0 for title slide, 1 for title and content

content_slide_layout = prs.slide_layouts[1] # 1 for title and content



# Title Slide

slide_0 = prs.slides.add_slide(title_slide_layout)

title = slide_0.shapes.title

subtitle = slide_0.placeholders[1]



title.text = "Supervised and Unsupervised Learning"

subtitle.text = "Understanding Key Concepts in Machine Learning\nPresenter Name\nDate"



# Introduction Slide

slide_1 = prs.slides.add_slide(content_slide_layout)

title = slide_1.shapes.title

content = slide_1.placeholders[1]



title.text = "Introduction to Machine Learning"

content.text = ("- Brief overview of machine learning\n"

                "- Types of machine learning\n"

                "  - Supervised Learning\n"

                "  - Unsupervised Learning\n"

                "  - Semi-Supervised Learning\n"

                "  - Reinforcement Learning")



# Supervised Learning Slide

slide_2 = prs.slides.add_slide(content_slide_layout)

title = slide_2.shapes.title

content = slide_2.placeholders[1]



title.text = "What is Supervised Learning?"

content.text = ("- Definition\n"

                "- Key characteristics\n"

                "- Use of labeled data\n"

                "- Common algorithms\n"

                "  - Linear Regression\n"

                "  - Logistic Regression\n"

                "  - Naive Bayes\n"

                "  - K-Nearest Neighbors (KNN)\n"

                "  - Random Forest")



# Applications of Supervised Learning Slide

slide_3 = prs.slides.add_slide(content_slide_layout)

title = slide_3.shapes.title

content = slide_3.placeholders[1]



title.text = "Applications of Supervised Learning"

content.text = ("- Predictive modeling\n"

                "- Classification problems\n"

                "- Real-world examples\n"

                "  - Spam detection\n"

                "  - Sentiment analysis\n"

                "  - Predictive maintenance")



# Unsupervised Learning Slide

slide_4 = prs.slides.add_slide(content_slide_layout)

title = slide_4.shapes.title

content = slide_4.placeholders[1]



title.text = "What is Unsupervised Learning?"

content.text = ("- Definition\n"

                "- Analyzing and clustering unlabeled data\n"

                "- Discovery of hidden patterns or data groupings")



# Key Tasks in Unsupervised Learning Slide

slide_5 = prs.slides.add_slide(content_slide_layout)

title = slide_5.shapes.title

content = slide_5.placeholders[1]



title.text = "Key Tasks in Unsupervised Learning"

content.text = ("- Clustering\n"

                "  - Exclusive (Hard) Clustering\n"

                "  - Overlapping (Soft) Clustering\n"

                "- Association Rules\n"

                "- Dimensionality Reduction")



# Clustering in Unsupervised Learning Slide

slide_6 = prs.slides.add_slide(content_slide_layout)

title = slide_6.shapes.title

content = slide_6.placeholders[1]



title.text = "Clustering in Unsupervised Learning"

content.text = ("- Definition and techniques\n"

                "- Exclusive Clustering (K-means)\n"

                "- Overlapping Clustering (Fuzzy K-means)")



# Hierarchical Clustering Slide

slide_7 = prs.slides.add_slide(content_slide_layout)

title = slide_7.shapes.title

content = slide_7.placeholders[1]



title.text = "Hierarchical Clustering"

content.text = ("- Agglomerative vs Divisive\n"

                "- Methods for measuring similarity\n"

                "  - Ward’s linkage\n"

                "  - Average linkage\n"

                "  - Complete linkage\n"

                "  - Single linkage")



# Probabilistic Clustering Slide

slide_8 = prs.slides.add_slide(content_slide_layout)

title = slide_8.shapes.title

content = slide_8.placeholders[1]



title.text = "Probabilistic Clustering"

content.text = ("- Definition\n"

                "- Gaussian Mixture Models (GMM)\n"

                "- Expectation-Maximization (EM) algorithm")



# Association Rules Slide

slide_9 = prs.slides.add_slide(content_slide_layout)

title = slide_9.shapes.title

content = slide_9.placeholders[1]



title.text = "Association Rules"

content.text = ("- Definition and use cases\n"

                "- Market basket analysis\n"

                "- Apriori algorithm")



# Dimensionality Reduction Slide

slide_10 = prs.slides.add_slide(content_slide_layout)

title = slide_10.shapes.title

content = slide_10.placeholders[1]



title.text = "Dimensionality Reduction"

content.text = ("- Importance in machine learning\n"

                "- Techniques\n"

                "  - Principal Component Analysis (PCA)\n"

                "  - Singular Value Decomposition (SVD)\n"

                "  - Autoencoders")



# Applications of Unsupervised Learning Slide

slide_11 = prs.slides.add_slide(content_slide_layout)

title = slide_11.shapes.title

content = slide_11.placeholders[1]



title.text = "Applications of Unsupervised Learning"

content.text = ("- Google News categorization\n"

                "- Computer vision\n"

                "- Medical imaging\n"

                "- Anomaly detection\n"

                "- Customer personas\n"

                "- Recommendation engines")



# Comparing Supervised and Unsupervised Learning Slide

slide_12 = prs.slides.add_slide(content_slide_layout)

title = slide_12.shapes.title

content = slide_12.placeholders[1]



title.text = "Comparing Supervised and Unsupervised Learning"

content.text = ("- Key differences\n"

                "- Advantages and disadvantages\n"

                "- Examples of use cases")



# Challenges of Unsupervised Learning Slide

slide_13 = prs.slides.add_slide(content_slide_layout)

title = slide_13.shapes.title

content = slide_13.placeholders[1]



title.text = "Challenges of Unsupervised Learning"

content.text = ("- Computational complexity\n"

                "- Training times\n"

                "- Risk of inaccurate results\n"

                "- Need for human intervention\n"

                "- Lack of transparency")



# Conclusion Slide

slide_14 = prs.slides.add_slide(content_slide_layout)

title = slide_14.shapes.title

content = slide_14.placeholders[1]



title.text = "Conclusion"

content.text = ("- Recap of supervised vs unsupervised learning\n"

                "- Importance of understanding both methods\n"

                "- Q&A")



# Save the presentation

pptx_path = "Supervised_vs_Unsupervised_Learning.pptx"

prs.save(pptx_path)



pptx_path